"""Generate project summary presentation using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Color palette
NAVY = RGBColor(0x0F, 0x1B, 0x33)
DARK_BG = RGBColor(0x16, 0x23, 0x44)
MID_BG = RGBColor(0x1A, 0x30, 0x50)
TEAL = RGBColor(0x2E, 0xC4, 0xB6)
CORAL = RGBColor(0xE7, 0x6F, 0x51)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x99, 0xAA)
LIGHT_GRAY = RGBColor(0xB0, 0xC0, 0xD0)
DARK_TEAL = RGBColor(0x0A, 0x47, 0x3F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, font_name="Arial", alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_accent_bar(slide, left, top, width, color):
    rect = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(0.05))  # 1 = rectangle
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()


def add_card(slide, left, top, width, height, title, body,
             title_color=TEAL, body_color=LIGHT_GRAY, border_color=TEAL):
    card = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_BG
    card.line.color.rgb = border_color
    card.line.width = Pt(2)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Arial"
    for line in body:
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(10)
        p2.font.color.rgb = body_color
        p2.font.name = "Arial"
        p2.space_before = Pt(4)


def add_section_header(slide, text):
    add_accent_bar(slide, 0.8, 0.5, 1.2, TEAL)
    add_text_box(slide, 0.8, 0.6, 10, 0.6, text, font_size=28,
                 color=WHITE, bold=True)


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, NAVY)
add_text_box(slide, 1.0, 2.0, 11, 0.6, "车载告警音分类系统", font_size=42,
             color=WHITE, bold=True)
add_accent_bar(slide, 1.0, 2.8, 2.0, CORAL)
add_text_box(slide, 1.0, 3.1, 11, 0.5,
             "Alarm Sound Classification — CNN + GRU Multi-Label Detection", font_size=20,
             color=TEAL)
add_text_box(slide, 1.0, 3.7, 11, 0.4,
             "端到端解决方案：从数据管线到推理部署", font_size=14, color=GRAY)

# ============================================================
# SLIDE 2: Project Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_section_header(slide, "项目概述")

add_card(slide, 0.8, 1.4, 3.8, 2.0, "输入",
         ["车载麦克风单声道音频（10-30s）",
          "可能包含：多种告警叠加、TTS 语音、发动机/行驶噪声"],
         title_color=TEAL, border_color=TEAL)
add_card(slide, 4.9, 1.4, 3.8, 2.0, "输出",
         ["检测到的告警类别及名称",
          "每类告警的响铃时间点（起止时间）",
          "响铃总次数（连续模式按周期推算）"],
         title_color=TEAL, border_color=TEAL)
add_card(slide, 9.0, 1.4, 3.8, 2.0, "场景",
         ["实时车载：单告警 / 多告警重叠 / +TTS",
          "不同距离、不同音量的识别鲁棒性"],
         title_color=TEAL, border_color=TEAL)
add_card(slide, 0.8, 3.7, 12.0, 1.0, "核心目标：从 6 类纯净单周期模板出发，仅用少量手动录制音频，训练一个能处理重叠告警、区分节奏差异、抗噪声干扰的鲁棒多标签分类器",
         [], title_color=CORAL, border_color=CORAL)

# ============================================================
# SLIDE 3: Alarm Types
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_section_header(slide, "6 类告警音 — 周期结构与特征")

alarms = [
    ("ACC", "周期 750ms", "3×250ms 脉冲（125ms响/125ms停）"),
    ("Caution1", "周期 500ms", "500ms 持续高能量纯音"),
    ("Caution2", "周期 1000ms", "500ms 响 + 500ms 停"),
    ("Dooropen", "周期 750ms", "500ms 高位 + 250ms 低位"),
    ("FCW", "周期 750ms", "3×250ms 脉冲（与 ACC 同节奏，不同频率）"),
    ("FEB", "周期 720ms", "3×240ms 脉冲（120ms响/120ms停）"),
]
for i, (name, dur, pat) in enumerate(alarms):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.4 + row * 2.5
    add_card(slide, x, y, 3.8, 2.2, name,
             [dur, pat], title_color=CORAL, border_color=DARK_TEAL)

add_text_box(slide, 0.8, 6.5, 12, 0.5,
             "关键区分点：ACC/FCW/FEB 都是脉冲型，靠频率区分。Caution1/Caution2 响的部分相同，靠'是否有停'区分。纯频谱不够，必须时序建模。",
             font_size=11, color=GRAY)

# ============================================================
# SLIDE 4: Core Challenges
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_section_header(slide, "核心挑战与解决方案")

challenges = [
    ("极少量训练数据", "每个告警仅 1 个纯净单周期文件，无大规模标注数据集",
     "模板铺放增强 → 完整周期作为原子，随机位移+噪声合成无限样本"),
    ("多告警重叠", "2-3 个告警同时响，频谱混合难以分离",
     "多热标签 + BCEWithLogitsLoss → 每个类独立 sigmoid，允许同时输出多个 1"),
    ("人声/TTS 干扰", "车内语音播报可能被误判为告警",
     "TTS 录音放入噪声池作为负样本；告警+TTS 叠加放入素材库"),
    ("节奏差异 vs 频谱相似", "ACC/FCW 同为脉冲型；Caution1/Caution2 响的部分频谱相同",
     "CRNN 架构：CNN 提取频率特征 + GRU 建模时序，学到\"响-停-响\"节奏"),
    ("合成训练 ≠ 真实推理", "合成数据训练的模型在真实录音上水土不服",
     "data/real_val/ 真实验证集 + 每个 epoch 自动评估"),
    ("未知重复次数+顺序", "告警种类、次数、顺序均不确定",
     "滑窗+概率峰检测+连续模式自动切换；连续告警按周期推算次数"),
]
for i, (title, desc, sol) in enumerate(challenges):
    y = 1.4 + (i % 3) * 1.9
    x = 0.8 + (i // 3) * 6.4
    add_card(slide, x, y, 6.0, 1.7, title,
             [desc, sol], title_color=CORAL, border_color=CORAL)

# ============================================================
# SLIDE 5: Data Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_section_header(slide, "数据管线")

add_card(slide, 0.8, 1.3, 5.8, 3.6, "训练数据流", [
    "data/alarms/ — 子目录=类名，文件名任意，支持单周期+长录音混放",
    "data/noise/ — TTS、发动机、人声、行驶噪声、合成噪声",
    "data_loader.py — 自动扫描目录、重采样22kHz、归一化",
    "augment.py — 模板铺放法：完整告警周期随机位移铺入1.2s窗口",
    "features.py — Mel频谱: 128频带×54帧",
    "dataset.py — PyTorch Dataset，每次__getitem__调用augment",
], title_color=TEAL, border_color=TEAL)

add_card(slide, 7.0, 1.3, 5.8, 3.6, "增强策略 & 真实验证", [
    "单告警 60%: 完整周期随机位移，加噪声",
    "双告警 20%: 两个周期不同位置，可能重叠",
    "纯噪声 20%: 从噪声池随机截取",
    "真实验证集 data/real_val/ — 不参与训练，每epoch自动评估",
    "labels.txt: 文件名=类1,类2,... （纯噪声右边留空）",
], title_color=CORAL, border_color=CORAL)

# ============================================================
# SLIDE 6: Model Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_section_header(slide, "模型架构: CRNN（CNN + Bi-GRU）")

arch_layers = [
    ("Mel 频谱输入", "1×128×54 (1.2s 音频)", TEAL),
    ("Conv1 → 32ch, 3×3, BN, ReLU, 2× Freq Pool", "32×64×54", TEAL),
    ("Conv2 → 64ch, 3×3, BN, ReLU, 2× Freq Pool", "64×32×54", TEAL),
    ("Conv3 → 128ch, 3×3, BN, ReLU, 2× Freq Pool", "128×16×54", TEAL),
    ("Conv4 → 128ch, 3×3, BN, ReLU (no pool)", "128×16×54", TEAL),
    ("AdaptiveAvgPool2d((1,None)) + Reshape", "54 timesteps × 128 features", TEAL),
    ("Bi-GRU ×2 (128 hidden, bidirectional, dropout 0.3)", "54 × 256", CORAL),
    ("MaxPool over Time", "256-dim", TEAL),
    ("Classifier: 256→64→8 + Dropout 0.3", "8-class logits", TEAL),
    ("Sigmoid → multi-hot output", "8 independent probabilities", CORAL),
]

for i, (desc, shape, color) in enumerate(arch_layers):
    y = 1.3 + i * 0.55
    add_text_box(slide, 0.8, y, 7.0, 0.5, f"{desc}", font_size=10,
                 color=color, bold=(color == CORAL))
    add_text_box(slide, 8.2, y, 4.5, 0.5, f"[{shape}]", font_size=9, color=GRAY)

add_text_box(slide, 0.8, 6.8, 12, 0.4,
             "总参数: 752,168  |  Loss: BCEWithLogitsLoss  |  Optimizer: Adam (LR=1e-3)  |  Scheduler: ReduceLROnPlateau",
             font_size=10, color=GRAY)

# ============================================================
# SLIDE 7: Architecture Details (why each layer)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_section_header(slide, "每一层做了什么 & 为什么这样设计")

details = [
    ("Conv1-3（频率池化，保留时间维）",
     "3层卷积逐步压缩频率(128→64→32→16)，时间维始终保留。通道数1→32→64→128 逐步增加特征容量。",
     "只压频率不压时间，为GRU提供完整54帧时序序列"),
    ("Conv4（无池化）",
     "128ch, 3×3卷积, 无池化。输出128×16×54。",
     "最后一层精炼特征，不再降维"),
    ("AdaptiveAvgPool2d((1,None))",
     "频率维全局平均池化到1。输出128×1×54→Reshape为54×128。",
     "将每帧的频域信息压缩为128维特征向量"),
    ("Bi-GRU ×2（核心时序层）",
     "54步×128维输入。128隐藏单元，2层，双向→输出256维/步。Dropout 0.3。",
     "前向GRU学'响→停'关系，后向GRU学'停→响'关系。双向拼接=完整上下文。2层=更深层次时序表征"),
    ("MaxPool over Time",
     "对GRU输出的54步取最大值。54×256→256。",
     "取最大激活值=确保任一帧触发就能响应。比Avg更适合稀疏事件检测"),
    ("独立Sigmoid输出",
     "8个Sigmoid独立输出，非Softmax。多热标签。",
     "窗口内可同时输出ACC=1+FCW=1，天然支持重叠告警"),
]

for i, (title, desc, why) in enumerate(details):
    y = 1.2 + i * 0.95
    add_card(slide, 0.8, y, 12.0, 0.85, title,
             [desc, f"原因：{why}"], title_color=TEAL, border_color=DARK_TEAL)

# ============================================================
# SLIDE 8: Training & Optimization
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_section_header(slide, "训练策略 & 调优历程")

add_card(slide, 0.8, 1.2, 6.0, 2.8, "数据增强比例演变", [
    "初始：NO_ALARM=15% SINGLE=30% DUAL=35% TRIPLE=20%",
    "问题：85%样本有告警，模型过度自信；25%为三告警，训练复杂",
    "优化后：NO_ALARM=20% SINGLE=60% DUAL=20% TRIPLE=0%",
    "新增：Caution1 pos_weight=1.5 防止被负样本淹没",
], title_color=CORAL, border_color=CORAL)

add_card(slide, 7.2, 1.2, 5.6, 2.8, "推理阈值扫描（test_01.wav）", [
    "0.5→4Caution1误报+5FCW误报+1FEB误报 (9总)",
    "0.6→3Caution2+0误报 ✓ (最佳)",
    "0.7-0.8→2Caution2+5FCW高置信度误报 (阈值杀不掉)",
    "0.9→2Caution2+3FCW (FCW conf=0.88-1.0, 训练层面问题)",
], title_color=TEAL, border_color=TEAL)

add_card(slide, 0.8, 4.3, 12.0, 2.8, "5 个关键突破", [
    "1. CNN→CRNN: GRU加入时序建模，Caution2节奏识别大幅提升",
    "2. 模板铺放法: 抛弃随机裁剪，用完整周期做原子。Caution2 F1 0.59→0.89",
    "3. 负样本调优: 30%→20%+Caution1加权，FCW/FEB误报消除",
    "4. 真实验证集: real_val 发现合成验证虚高问题，5/7→6/7提升可追踪",
    "5. _is_continuous修复: 全局比例→最长连续段≥3周期，FEB 1次→12次，ACC 1次→10次",
], title_color=CORAL, border_color=CORAL)

# ============================================================
# SLIDE 9: Inference Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_section_header(slide, "推理管线: 从音频到响铃时间点")

inf_steps = [
    ("1. 加载模型", "自动检测CNN/CRNN类型，加载alarm_names，CPU推理"),
    ("2. 滑动窗口", "1.2s窗口 / 0.1s步长。每窗口→mel→模型→8类概率向量"),
    ("3. 概率序列", "每类prob(t)，0.1s分辨率，连续概率曲线"),
    ("4. 模式判断", "连续？最长活跃段≥3周期。离散？find_peaks找峰。"),
    ("5. 响铃输出", "离散：每峰=1次响。连续：按cycle_duration步进。含时间+置信度+数量。"),
]
for i, (title, desc) in enumerate(inf_steps):
    x = 0.8 + i * 2.5
    add_card(slide, x, 1.3, 2.3, 2.5, title, [desc], title_color=TEAL, border_color=DARK_TEAL)

add_card(slide, 0.8, 4.1, 5.8, 2.8, "连续模式 (detect_continuous_rings)", [
    "告警持续响(Dooropen连响7s)→概率平台",
    "从首次检测开始，按cycle_duration步进",
    "输出近似数量(±1-2次，受窗口拖尾影响)",
    "标记为(approx, continuous)",
], title_color=TEAL, border_color=TEAL)
add_card(slide, 7.0, 4.1, 5.8, 2.8, "离散模式 + 可调参数", [
    "告警断续响→概率峰。find_peaks(height, distance, prominence)",
    "--threshold 置信度门槛(默认0.6)",
    "--min-ring-spacing 最小峰间距ms(默认300)",
    "--model 指定模型路径",
], title_color=CORAL, border_color=CORAL)

# ============================================================
# SLIDE 10: Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_section_header(slide, "最终成果 & 项目文件结构")

add_card(slide, 0.8, 1.2, 3.8, 2.5, "模型指标 (合成验证)", [
    "Macro F1: 0.92-0.96",
    "Caution1: 0.975  Caution2: 0.83",
    "ACC: 0.95  Dooropen: 0.93",
    "FCW: 0.96  FEB: 0.99",
], title_color=TEAL, border_color=TEAL)
add_card(slide, 4.9, 1.2, 3.8, 2.5, "真实验证 (real_val)", [
    "6/7 真实文件识别通过",
    "FEB 1→12次，ACC 1→10次",
    "(连续模式修复后)",
    "Caution2在test_01正确检测2次",
], title_color=CORAL, border_color=CORAL)
add_card(slide, 9.0, 1.2, 3.8, 2.5, "推理性能", [
    "CPU: 10s音频≈5-10s推理",
    "时间分辨率: 0.1s",
    "连续/离散双模式自动切换",
    "阈值0.6最佳",
], title_color=TEAL, border_color=TEAL)

add_card(slide, 0.8, 4.0, 12.0, 2.8, "项目文件结构 & 技术栈", [
    "src/: config / data_loader / augment / dataset / model / features / train / infer / record_classify",
    "data/: alarms(子目录=类名) / noise(背景噪声) / real_val(真实验证集+labels.txt)",
    "models/: 训练输出按run_YYYYMMDD_HHMMSS存放, best_model.pt为最新符号链接",
    "技术栈: PyTorch + Librosa + SoundFile + scikit-learn + SciPy(find_peaks)",
    "架构: CNN骨干(4层Conv) + Bi-GRU(2层×128) + 分类头(256→64→8) + 多热标签 + BCEWithLogitsLoss",
    "增强: 模板铺放法 + Mel频谱(128×54) + 随机SNR噪声(-5~+20dB) + 多告警重叠合成",
], title_color=TEAL, border_color=TEAL)

# ============================================================
# SLIDE 11: Sequence / Flow Diagram
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, NAVY)
add_section_header(slide, "端到端时序流程")

# --- Training Pipeline Row ---
add_text_box(slide, 0.8, 1.3, 2.5, 0.4, "训练管线",
             font_size=14, color=CORAL, bold=True)

train_steps = [
    ("data/alarms/\ndata/noise/", TEAL),
    ("data_loader.py\n扫描/重采样/归一化", TEAL),
    ("augment.py\n模板铺放+噪声混合", TEAL),
    ("features.py\nMel频谱(128×54)", TEAL),
    ("CRNN\n前向传播", CORAL),
    ("BCEWithLogits\nLoss + 反向传播", CORAL),
    ("ReduceLROnPlateau\n早停/保存最佳模型", TEAL),
]

for i, (text, color) in enumerate(train_steps):
    x = 0.5 + i * 1.85
    box = slide.shapes.add_shape(1, Inches(x), Inches(1.8), Inches(1.7), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = DARK_BG
    box.line.color.rgb = color; box.line.width = Pt(1.5)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(7.5); p.font.color.rgb = color
    p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER
    # Arrow to next
    if i < len(train_steps) - 1:
        add_text_box(slide, x + 1.7, 2.0, 0.2, 0.4, "→",
                     font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# Iteration arrow
add_text_box(slide, 0.8, 2.7, 12, 0.4,
             "每 epoch 4000 个合成样本 → 验证(合成1000+真实验证) → 重复 EPOCHS 轮",
             font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

# --- Inference Pipeline Row ---
add_text_box(slide, 0.8, 3.2, 2.5, 0.4, "推理管线",
             font_size=14, color=TEAL, bold=True)

inf_steps = [
    ("test.wav\n原始音频", CORAL),
    ("librosa.load\n重采样22kHz", TEAL),
    ("滑动窗口\n1.2s窗/0.1s步长", TEAL),
    ("features.py\nMel频谱(128×54)", TEAL),
    ("CRNN\n前向传播", CORAL),
    ("Sigmoid→概率\n8类prob(t)曲线", TEAL),
    ("find_peaks或\n连续模式判断", CORAL),
    ("输出：类别+时间\n+次数+置信度", CORAL),
]

for i, (text, color) in enumerate(inf_steps):
    x = 0.2 + i * 1.6 if i < 7 else 0.2 + i * 1.6
    x = 0.2 + i * 1.62
    box = slide.shapes.add_shape(1, Inches(x), Inches(3.7), Inches(1.5), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = DARK_BG
    box.line.color.rgb = color; box.line.width = Pt(1.5)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(7.5); p.font.color.rgb = color
    p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER
    if i < len(inf_steps) - 1:
        add_text_box(slide, x + 1.5, 3.9, 0.15, 0.4, "→",
                     font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# --- Bottom: Feedback Loop ---
add_card(slide, 0.8, 4.8, 5.8, 2.0, "训练反馈闭环",
         ["val_macro_f1 监控合成验证是否在学",
          "real_val 真实验证防止合成过拟合",
          "早停 patience=12 避免无用训练",
          "每 epoch 后 real-val 结果对比 → 决定何时停止"],
         title_color=CORAL, border_color=CORAL)

add_card(slide, 7.0, 4.8, 5.8, 2.0, "推理后处理（infer.py）",
         ["_is_continuous: 最长连续活跃 ≥ 3×周期 → 连续模式",
          "_detect_continuous_rings: 按步长 → 估算次数",
          "_detect_discrete_rings: find_peaks → 每峰=1次",
          "输出：每类告警 × 几次 × 在哪些时间点 × 置信度"],
         title_color=TEAL, border_color=TEAL)

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "Alarm_Sound_Classification_Project.pptx")
prs.save(output_path)
print(f"Saved to: {output_path}")
